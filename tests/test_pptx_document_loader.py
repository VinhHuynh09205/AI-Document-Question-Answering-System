import time
from pathlib import Path

from pptx import Presentation

from app.services.document_loaders.pptx_document_loader import PptxDocumentLoader


def test_pptx_loader_rejects_provider_marker_noise_note() -> None:
    assert (
        PptxDocumentLoader._is_useful_image_note(
            "slide image 1 local_ocr extracted text",
            slide_prefers_cjk=False,
        )
        is False
    )


def test_pptx_loader_rejects_latin_ocr_when_slide_is_cjk() -> None:
    assert (
        PptxDocumentLoader._is_useful_image_note(
            "photpah olibrany ehooln",
            slide_prefers_cjk=True,
        )
        is False
    )


def test_pptx_loader_accepts_meaningful_latin_note() -> None:
    assert (
        PptxDocumentLoader._is_useful_image_note(
            "Customer retention increased by 12 percent year over year.",
            slide_prefers_cjk=False,
        )
        is True
    )


def test_pptx_loader_detects_duplicate_note_from_slide_text() -> None:
    assert (
        PptxDocumentLoader._is_duplicate_image_note(
            "Muc tieu chinh",
            "Muc tieu chinh cua bai hoc la nang cao ky nang.",
            [],
        )
        is True
    )


def test_pptx_loader_accepts_numeric_chart_note_without_two_latin_words() -> None:
    assert (
        PptxDocumentLoader._is_useful_image_note(
            "trieungoi 26.9 24.2 22.3 20.8",
            slide_prefers_cjk=False,
        )
        is True
    )


def test_pptx_loader_skips_dense_text_slide_with_single_picture() -> None:
    loader = PptxDocumentLoader(
        image_understanding_service=object(),
        text_char_threshold_for_image_analysis=900,
    )

    should_analyze = loader._should_analyze_slide_images(
        text_snapshot="word " * 220,
        total_images_analyzed=0,
        analyzed_slides=0,
        started_at=time.perf_counter(),
        picture_count=1,
    )

    assert should_analyze is False


def test_pptx_loader_allows_image_heavy_slide_with_limited_text() -> None:
    loader = PptxDocumentLoader(
        image_understanding_service=object(),
        text_char_threshold_for_image_analysis=900,
    )

    should_analyze = loader._should_analyze_slide_images(
        text_snapshot="Key chart summary",
        total_images_analyzed=0,
        analyzed_slides=0,
        started_at=time.perf_counter(),
        picture_count=3,
    )

    assert should_analyze is True


def test_pptx_loader_honors_max_slides_with_image_analysis() -> None:
    loader = PptxDocumentLoader(
        image_understanding_service=object(),
        max_slides_with_image_analysis=2,
    )

    should_analyze = loader._should_analyze_slide_images(
        text_snapshot="small text",
        total_images_analyzed=0,
        analyzed_slides=2,
        started_at=time.perf_counter(),
        picture_count=2,
    )

    assert should_analyze is False


def test_pptx_loader_emits_structured_slide_blocks_and_layout(tmp_path: Path) -> None:
    pptx_path = tmp_path / "demo.pptx"

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Roadmap"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Milestone planning"
    bullet = body.add_paragraph()
    bullet.text = "Phase 1 deliverables"
    bullet.level = 1
    slide.notes_slide.notes_text_frame.text = "Discuss risk mitigation options."

    table_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    table_shape = table_slide.shapes.add_table(2, 2, left=0, top=0, width=2_000_000, height=800_000)
    table_shape.table.cell(0, 0).text = "Metric"
    table_shape.table.cell(0, 1).text = "Value"
    table_shape.table.cell(1, 0).text = "Revenue"
    table_shape.table.cell(1, 1).text = "120"

    presentation.save(pptx_path)

    loader = PptxDocumentLoader()
    docs = loader.load(pptx_path)

    assert len(docs) == 2

    first_slide = docs[0]
    assert first_slide.metadata.get("content_type") == "slide_structured"
    assert first_slide.metadata.get("slide_number") == 1
    assert first_slide.metadata.get("slide_layout")
    assert first_slide.metadata.get("text_block_count", 0) >= 1
    assert first_slide.metadata.get("has_notes") is True

    blocks = first_slide.metadata.get("slide_blocks")
    assert isinstance(blocks, list)
    assert blocks
    assert all("reading_order" in block for block in blocks)
    assert all("position" in block for block in blocks)

    second_slide = docs[1]
    assert second_slide.metadata.get("slide_number") == 2
    assert second_slide.metadata.get("has_table") is True
