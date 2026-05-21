from __future__ import annotations

import hashlib
import logging
import re
import time
from pathlib import Path

from langchain_core.documents import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.interfaces.document_loader import IDocumentLoader
from app.services.interfaces.image_understanding_service import IImageUnderstandingService


logger = logging.getLogger(__name__)


class PptxDocumentLoader(IDocumentLoader):
    _OCR_MARKER_RE = re.compile(
        r"(local_ocr|local_vision|image\s*analysis|slide\s*image|provider[:=])",
        re.IGNORECASE,
    )
    _LATIN_WORD_RE = re.compile(r"[A-Za-z\u00C0-\u024F]{2,}")
    _CJK_RE = re.compile(r"[\u3040-\u30FF\u3400-\u4DBF\u4E00-\u9FFF]")
    _NUMERIC_TOKEN_RE = re.compile(r"\d+(?:[.,]\d+)?")
    _DECORATIVE_LINE_RE = re.compile(
        r"^(?:\d{1,3}|page\s*\d{1,3}|slide\s*\d{1,3}|\d{1,3}/\d{1,3})$",
        re.IGNORECASE,
    )

    def __init__(
        self,
        image_understanding_service: IImageUnderstandingService | None = None,
        *,
        max_images_per_slide: int = 2,
        max_images_per_document: int = 24,
        max_slides_with_image_analysis: int = 12,
        text_char_threshold_for_image_analysis: int = 950,
        max_image_analysis_seconds: float = 25.0,
    ) -> None:
        self._image_understanding_service = image_understanding_service
        self._max_images_per_slide = max(1, int(max_images_per_slide))
        self._max_images_per_document = max(1, int(max_images_per_document))
        self._max_slides_with_image_analysis = max(1, int(max_slides_with_image_analysis))
        self._text_char_threshold_for_image_analysis = max(250, int(text_char_threshold_for_image_analysis))
        self._max_image_analysis_seconds = max(1.0, float(max_image_analysis_seconds))

    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() == ".pptx"

    def load(self, file_path: Path) -> list[Document]:
        try:
            documents = self._load_primary(file_path)
            logger.info(
                "pptx_loader_primary_success file=%s documents=%s",
                file_path.name,
                len(documents),
            )
            return documents
        except Exception:
            logger.warning(
                "pptx_loader_primary_failed file=%s fallback=simple_slide_parser",
                file_path.name,
                exc_info=True,
            )

        try:
            documents = self._load_fallback(file_path)
            logger.info(
                "pptx_loader_fallback_success file=%s documents=%s",
                file_path.name,
                len(documents),
            )
            return documents
        except Exception as exc:
            logger.error("pptx_loader_fallback_failed file=%s", file_path.name, exc_info=True)
            raise RuntimeError(f"Failed to parse presentation {file_path.name}") from exc

    def _load_primary(self, file_path: Path) -> list[Document]:
        presentation = Presentation(str(file_path))
        repeated_lines = self._detect_repeated_slide_lines(presentation)

        image_analysis_started_at = time.perf_counter()
        total_images_analyzed = 0
        analyzed_slides = 0
        seen_image_hashes: set[str] = set()

        documents: list[Document] = []
        total_slides = 0
        title_count = 0
        text_blocks = 0
        table_blocks = 0
        chart_blocks = 0
        image_blocks = 0
        notes_blocks = 0
        ocr_blocks = 0
        vision_blocks = 0

        for slide_index, slide in enumerate(presentation.slides, start=1):
            total_slides += 1
            try:
                document, stats, image_stats = self._extract_slide_document(
                    file_path=file_path,
                    slide=slide,
                    slide_index=slide_index,
                    repeated_lines=repeated_lines,
                    image_analysis_started_at=image_analysis_started_at,
                    total_images_analyzed=total_images_analyzed,
                    analyzed_slides=analyzed_slides,
                    seen_image_hashes=seen_image_hashes,
                )
            except Exception:
                logger.warning(
                    "pptx_loader_slide_primary_failed file=%s slide=%s fallback=simple_slide",
                    file_path.name,
                    slide_index,
                    exc_info=True,
                )
                document = self._extract_slide_document_fallback(file_path=file_path, slide=slide, slide_index=slide_index)
                stats = {
                    "has_title": False,
                    "text_blocks": 0,
                    "table_blocks": 0,
                    "chart_blocks": 0,
                    "image_blocks": 0,
                    "notes_blocks": 0,
                    "ocr_blocks": 0,
                    "vision_blocks": 0,
                }
                image_stats = {
                    "total_images_analyzed": total_images_analyzed,
                    "analyzed_slides": analyzed_slides,
                }

            total_images_analyzed = int(image_stats.get("total_images_analyzed", total_images_analyzed))
            analyzed_slides = int(image_stats.get("analyzed_slides", analyzed_slides))

            if document is None:
                continue
            documents.append(document)

            title_count += 1 if bool(stats.get("has_title")) else 0
            text_blocks += int(stats.get("text_blocks", 0))
            table_blocks += int(stats.get("table_blocks", 0))
            chart_blocks += int(stats.get("chart_blocks", 0))
            image_blocks += int(stats.get("image_blocks", 0))
            notes_blocks += int(stats.get("notes_blocks", 0))
            ocr_blocks += int(stats.get("ocr_blocks", 0))
            vision_blocks += int(stats.get("vision_blocks", 0))

        if not documents:
            raise ValueError("No extractable slide content")

        logger.info(
            "pptx_loader_summary file=%s slides=%s titles=%s text_blocks=%s table_blocks=%s chart_blocks=%s image_blocks=%s notes_blocks=%s ocr_blocks=%s vision_blocks=%s images_analyzed=%s",
            file_path.name,
            total_slides,
            title_count,
            text_blocks,
            table_blocks,
            chart_blocks,
            image_blocks,
            notes_blocks,
            ocr_blocks,
            vision_blocks,
            total_images_analyzed,
        )
        return documents

    def _load_fallback(self, file_path: Path) -> list[Document]:
        presentation = Presentation(str(file_path))
        documents: list[Document] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            fallback_doc = self._extract_slide_document_fallback(file_path=file_path, slide=slide, slide_index=slide_index)
            if fallback_doc is None:
                continue
            documents.append(fallback_doc)

        if not documents:
            raise ValueError("No extractable slide content in fallback parser")
        return documents

    def _extract_slide_document(
        self,
        *,
        file_path: Path,
        slide,
        slide_index: int,
        repeated_lines: set[str],
        image_analysis_started_at: float,
        total_images_analyzed: int,
        analyzed_slides: int,
        seen_image_hashes: set[str],
    ) -> tuple[Document | None, dict[str, int | bool], dict[str, int]]:
        title_shape = getattr(slide.shapes, "title", None)
        slide_title = str(getattr(title_shape, "text", "") or "").strip()
        slide_layout = str(getattr(getattr(slide, "slide_layout", None), "name", "") or "").strip()

        ordered_shapes = self._ordered_shapes(slide.shapes)
        shape_order_map = {id(shape): order for order, shape in enumerate(ordered_shapes, start=1)}

        blocks: list[dict[str, object]] = []
        slide_text_snapshot_lines: list[str] = []
        picture_shapes: list = []

        text_block_count = 0
        table_block_count = 0
        chart_block_count = 0

        for shape in ordered_shapes:
            reading_order = int(shape_order_map.get(id(shape), 0) or 0)
            object_type = self._resolve_object_type(shape)
            position = self._shape_position(shape)

            if title_shape is not None and shape == title_shape:
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shapes.append(shape)
                continue

            text_lines, has_bullet = self._extract_text_lines(shape, repeated_lines=repeated_lines, slide_title=slide_title)
            if text_lines:
                text_block_count += 1
                block_type = "bullet" if has_bullet else "textbox"
                content = "\n".join(text_lines)
                blocks.append(
                    self._build_block(
                        reading_order=reading_order,
                        block_type=block_type,
                        object_type=object_type,
                        position=position,
                        content=content,
                    )
                )
                slide_text_snapshot_lines.extend(text_lines)

            if shape.has_table:
                table_lines = self._extract_table_lines(shape)
                if table_lines:
                    table_block_count += 1
                    blocks.append(
                        self._build_block(
                            reading_order=reading_order,
                            block_type="table",
                            object_type=object_type,
                            position=position,
                            content="\n".join(table_lines),
                        )
                    )
                    slide_text_snapshot_lines.extend(table_lines)

            if getattr(shape, "has_chart", False):
                chart_lines = self._extract_chart_lines(shape)
                if chart_lines:
                    chart_block_count += 1
                    blocks.append(
                        self._build_block(
                            reading_order=reading_order,
                            block_type="chart",
                            object_type=object_type,
                            position=position,
                            content="\n".join(chart_lines),
                        )
                    )
                    slide_text_snapshot_lines.extend(chart_lines)

        slide_text_snapshot = "\n".join([slide_title, *slide_text_snapshot_lines]).strip()
        slide_prefers_cjk = self._contains_substantial_cjk(slide_text_snapshot)

        should_analyze_images = self._should_analyze_slide_images(
            text_snapshot=slide_text_snapshot,
            total_images_analyzed=total_images_analyzed,
            analyzed_slides=analyzed_slides,
            started_at=image_analysis_started_at,
            picture_count=len(picture_shapes),
        )

        ocr_blocks = 0
        vision_blocks = 0
        image_block_count = 0

        if should_analyze_images and picture_shapes:
            analyzed_slides += 1

        images_seen_on_slide = 0
        for picture_shape in picture_shapes:
            reading_order = int(shape_order_map.get(id(picture_shape), 0) or 0)
            object_type = self._resolve_object_type(picture_shape)
            position = self._shape_position(picture_shape)

            image_block_count += 1
            blocks.append(
                self._build_block(
                    reading_order=reading_order,
                    block_type="image",
                    object_type=object_type,
                    position=position,
                    content="Image object detected on slide.",
                )
            )

            if not should_analyze_images:
                continue
            if images_seen_on_slide >= self._max_images_per_slide:
                continue
            if total_images_analyzed >= self._max_images_per_document:
                continue
            if (time.perf_counter() - image_analysis_started_at) >= self._max_image_analysis_seconds:
                continue

            image_bytes = getattr(getattr(picture_shape, "image", None), "blob", b"")
            if not image_bytes:
                continue

            image_hash = self._fingerprint_image_bytes(image_bytes)
            if image_hash in seen_image_hashes:
                continue
            seen_image_hashes.add(image_hash)

            images_seen_on_slide += 1
            result = self._image_understanding_service.analyze_image(
                image_bytes,
                source=str(file_path),
                hint=f"pptx slide {slide_index} image {images_seen_on_slide}",
            )
            total_images_analyzed += 1

            provider_name = str(result.provider or "").strip().lower()
            note_text = self._normalize_image_note_text(str(result.text or ""))
            if not note_text:
                continue
            if not self._is_useful_image_note(note_text, slide_prefers_cjk=slide_prefers_cjk):
                continue
            if self._is_duplicate_image_note(note_text, slide_text_snapshot, [str(block.get("content") or "") for block in blocks]):
                continue

            if "ocr" in provider_name:
                ocr_blocks += 1
                block_type = "image_ocr"
            else:
                vision_blocks += 1
                block_type = "image_vision"

            blocks.append(
                self._build_block(
                    reading_order=reading_order,
                    block_type=block_type,
                    object_type=object_type,
                    position=position,
                    content=note_text,
                )
            )

        notes_text = ""
        if getattr(slide, "has_notes_slide", False):
            try:
                notes_text = self._normalize_slide_line(str(slide.notes_slide.notes_text_frame.text or ""))
            except Exception:
                notes_text = ""

        notes_block_count = 0
        if notes_text:
            notes_block_count = 1
            blocks.append(
                self._build_block(
                    reading_order=len(ordered_shapes) + 1,
                    block_type="speaker_notes",
                    object_type="notes",
                    position="",
                    content=notes_text,
                )
            )

        blocks = [block for block in blocks if str(block.get("content") or "").strip()]
        blocks.sort(key=lambda item: int(item.get("reading_order", 0) or 0))

        if not blocks and not slide_title:
            return None, {}, {
                "total_images_analyzed": total_images_analyzed,
                "analyzed_slides": analyzed_slides,
            }

        content = self._build_slide_document_content(
            file_name=file_path.name,
            slide_index=slide_index,
            slide_title=slide_title,
            slide_layout=slide_layout,
            blocks=blocks,
        )

        document = Document(
            page_content=content,
            metadata={
                "source": str(file_path),
                "extension": ".pptx",
                "content_type": "slide_structured",
                "file_name": file_path.name,
                "slide": slide_index,
                "slide_number": slide_index,
                "slide_title": slide_title,
                "slide_layout": slide_layout,
                "section_title": slide_title or f"Slide {slide_index}",
                "has_table": table_block_count > 0,
                "has_chart": chart_block_count > 0,
                "has_image": image_block_count > 0,
                "has_notes": bool(notes_text),
                "image_count": len(picture_shapes),
                "image_analysis_applied": bool(ocr_blocks or vision_blocks),
                "ocr_applied": bool(ocr_blocks > 0),
                "text_block_count": text_block_count,
                "table_block_count": table_block_count,
                "chart_block_count": chart_block_count,
                "image_block_count": image_block_count,
                "notes_block_count": notes_block_count,
                "slide_blocks": blocks,
            },
        )

        stats: dict[str, int | bool] = {
            "has_title": bool(slide_title),
            "text_blocks": text_block_count,
            "table_blocks": table_block_count,
            "chart_blocks": chart_block_count,
            "image_blocks": image_block_count,
            "notes_blocks": notes_block_count,
            "ocr_blocks": ocr_blocks,
            "vision_blocks": vision_blocks,
        }
        image_stats = {
            "total_images_analyzed": total_images_analyzed,
            "analyzed_slides": analyzed_slides,
        }
        return document, stats, image_stats

    def _extract_slide_document_fallback(self, *, file_path: Path, slide, slide_index: int) -> Document | None:
        title_shape = getattr(slide.shapes, "title", None)
        slide_title = str(getattr(title_shape, "text", "") or "").strip()
        slide_layout = str(getattr(getattr(slide, "slide_layout", None), "name", "") or "").strip()

        lines: list[str] = []
        for shape in slide.shapes:
            if title_shape is not None and shape == title_shape:
                continue
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = self._normalize_slide_line(paragraph.text)
                if text:
                    lines.append(text)

        notes_text = ""
        if getattr(slide, "has_notes_slide", False):
            try:
                notes_text = self._normalize_slide_line(str(slide.notes_slide.notes_text_frame.text or ""))
            except Exception:
                notes_text = ""

        blocks: list[dict[str, object]] = []
        if lines:
            blocks.append(
                self._build_block(
                    reading_order=1,
                    block_type="textbox",
                    object_type="text",
                    position="",
                    content="\n".join(lines[:80]),
                )
            )
        if notes_text:
            blocks.append(
                self._build_block(
                    reading_order=2,
                    block_type="speaker_notes",
                    object_type="notes",
                    position="",
                    content=notes_text,
                )
            )

        if not blocks and not slide_title:
            return None

        content = self._build_slide_document_content(
            file_name=file_path.name,
            slide_index=slide_index,
            slide_title=slide_title,
            slide_layout=slide_layout,
            blocks=blocks,
        )

        return Document(
            page_content=content,
            metadata={
                "source": str(file_path),
                "extension": ".pptx",
                "content_type": "slide_structured",
                "file_name": file_path.name,
                "slide": slide_index,
                "slide_number": slide_index,
                "slide_title": slide_title,
                "slide_layout": slide_layout,
                "section_title": slide_title or f"Slide {slide_index}",
                "has_table": False,
                "has_chart": False,
                "has_image": False,
                "has_notes": bool(notes_text),
                "image_count": 0,
                "image_analysis_applied": False,
                "ocr_applied": False,
                "text_block_count": 1 if lines else 0,
                "table_block_count": 0,
                "chart_block_count": 0,
                "image_block_count": 0,
                "notes_block_count": 1 if notes_text else 0,
                "slide_blocks": blocks,
            },
        )

    @classmethod
    def _ordered_shapes(cls, shapes) -> list:
        return sorted(
            list(shapes),
            key=lambda shape: (
                int(getattr(shape, "top", 0) or 0),
                int(getattr(shape, "left", 0) or 0),
            ),
        )

    @staticmethod
    def _resolve_object_type(shape) -> str:
        shape_type = getattr(shape, "shape_type", None)
        if shape_type is None:
            return "unknown"
        name = getattr(shape_type, "name", "")
        if name:
            return str(name).lower()
        return str(shape_type).lower()

    @staticmethod
    def _shape_position(shape) -> str:
        left = int(getattr(shape, "left", 0) or 0)
        top = int(getattr(shape, "top", 0) or 0)
        width = int(getattr(shape, "width", 0) or 0)
        height = int(getattr(shape, "height", 0) or 0)
        return f"x={left},y={top},w={width},h={height}"

    def _extract_text_lines(self, shape, *, repeated_lines: set[str], slide_title: str) -> tuple[list[str], bool]:
        if not shape.has_text_frame:
            return [], False

        lines: list[str] = []
        has_bullet = False
        for paragraph in shape.text_frame.paragraphs:
            text = self._normalize_slide_line(paragraph.text)
            if not text:
                continue
            if text.casefold() in repeated_lines:
                continue
            if self._looks_decorative_slide_line(text):
                continue
            if slide_title and text.casefold() == slide_title.casefold():
                continue

            if int(getattr(paragraph, "level", 0) or 0) > 0:
                has_bullet = True
            lines.append(text)

        deduped: list[str] = []
        seen: set[str] = set()
        for line in lines:
            key = line.casefold()
            if key in seen:
                continue
            seen.add(key)
            deduped.append(line)

        return deduped[:40], has_bullet

    @staticmethod
    def _build_block(
        *,
        reading_order: int,
        block_type: str,
        object_type: str,
        position: str,
        content: str,
    ) -> dict[str, object]:
        return {
            "reading_order": reading_order,
            "block_type": str(block_type or "").strip(),
            "object_type": str(object_type or "").strip(),
            "position": str(position or "").strip(),
            "content": str(content or "").strip(),
        }

    @staticmethod
    def _build_slide_document_content(
        *,
        file_name: str,
        slide_index: int,
        slide_title: str,
        slide_layout: str,
        blocks: list[dict[str, object]],
    ) -> str:
        lines: list[str] = [
            f"File: {file_name}",
            f"Slide {slide_index}",
            f"Title: {slide_title}" if slide_title else "",
            "Content:",
        ]

        for block in blocks:
            block_type = str(block.get("block_type") or "object")
            content = str(block.get("content") or "").strip()
            if not content:
                continue

            if block_type in {"table", "chart", "image_ocr", "image_vision", "speaker_notes"}:
                label = {
                    "table": "Table",
                    "chart": "Chart",
                    "image_ocr": "Image OCR",
                    "image_vision": "Image description",
                    "speaker_notes": "Speaker notes",
                }.get(block_type, "Slide content")
                lines.append(f"- {label}:")

            for content_line in content.splitlines():
                compact = str(content_line).strip()
                if not compact:
                    continue
                lines.append(f"  {compact}")

        return "\n".join(line for line in lines if line).strip()

    @classmethod
    def _normalize_slide_line(cls, text: str) -> str:
        return re.sub(r"\s+", " ", str(text or "")).strip(" -•\t")

    @classmethod
    def _looks_decorative_slide_line(cls, text: str) -> bool:
        compact = str(text or "").strip()
        if not compact:
            return True
        if cls._DECORATIVE_LINE_RE.match(compact):
            return True
        if len(compact) <= 2:
            return True
        return False

    def _should_analyze_slide_images(
        self,
        *,
        text_snapshot: str,
        total_images_analyzed: int,
        analyzed_slides: int,
        started_at: float,
        picture_count: int,
    ) -> bool:
        if self._image_understanding_service is None:
            return False

        if picture_count <= 0:
            return False

        if total_images_analyzed >= self._max_images_per_document:
            return False

        if analyzed_slides >= self._max_slides_with_image_analysis:
            return False

        elapsed_seconds = time.perf_counter() - started_at
        if elapsed_seconds >= self._max_image_analysis_seconds:
            return False

        compact_text = str(text_snapshot or "").strip()
        if not compact_text:
            return True

        if len(compact_text) >= self._text_char_threshold_for_image_analysis:
            if picture_count < 2:
                return False
            if len(compact_text) >= int(self._text_char_threshold_for_image_analysis * 1.4):
                return False

        if picture_count >= 2 and len(compact_text) < int(self._text_char_threshold_for_image_analysis * 1.25):
            return True

        if len(compact_text) >= self._text_char_threshold_for_image_analysis:
            return False

        if len(self._LATIN_WORD_RE.findall(compact_text)) >= 120 and picture_count < 3:
            return False

        return True

    @staticmethod
    def _fingerprint_image_bytes(image_bytes: bytes) -> str:
        return hashlib.sha1(image_bytes).hexdigest()[:24]

    @classmethod
    def _detect_repeated_slide_lines(cls, presentation: Presentation) -> set[str]:
        line_counts: dict[str, int] = {}
        total_slides = max(1, len(presentation.slides))
        threshold = max(3, int(total_slides * 0.5))

        for slide in presentation.slides:
            seen_in_slide: set[str] = set()
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    line = cls._normalize_slide_line(paragraph.text)
                    if not line:
                        continue
                    if cls._looks_decorative_slide_line(line):
                        continue

                    key = line.casefold()
                    if key in seen_in_slide:
                        continue
                    seen_in_slide.add(key)
                    line_counts[key] = line_counts.get(key, 0) + 1

        return {
            line for line, count in line_counts.items()
            if count >= threshold and len(line) <= 90
        }

    @classmethod
    def _normalize_image_note_text(cls, text: str) -> str:
        normalized = str(text or "").replace("\r\n", "\n").replace("\r", "\n")
        cleaned_lines: list[str] = []
        seen: set[str] = set()

        for raw_line in normalized.splitlines():
            line = re.sub(r"\s+", " ", raw_line).strip(" -•\t")
            if not line:
                continue

            lowered = line.lower()
            if lowered in seen:
                continue

            seen.add(lowered)
            cleaned_lines.append(line)
            if len(cleaned_lines) >= 5:
                break

        return "\n".join(cleaned_lines).strip()

    @classmethod
    def _contains_substantial_cjk(cls, text: str) -> bool:
        return len(cls._CJK_RE.findall(str(text or ""))) >= 6

    @classmethod
    def _is_useful_image_note(
        cls,
        note_text: str,
        *,
        slide_prefers_cjk: bool,
    ) -> bool:
        compact = " ".join(str(note_text or "").split()).strip()
        if len(compact) < 8:
            return False

        if cls._OCR_MARKER_RE.search(compact):
            return False

        has_cjk = bool(cls._CJK_RE.search(compact))
        if slide_prefers_cjk and not has_cjk:
            return False

        stripped_symbols = re.sub(
            r"[A-Za-z\u00C0-\u024F0-9\u3040-\u30FF\u3400-\u4DBF\u4E00-\u9FFF\s]",
            "",
            compact,
        )
        if len(stripped_symbols) > len(compact) * 0.35:
            return False

        if not has_cjk:
            latin_word_count = len(cls._LATIN_WORD_RE.findall(compact))
            if latin_word_count < 2:
                numeric_count = len(cls._NUMERIC_TOKEN_RE.findall(compact))
                if not (
                    (latin_word_count >= 1 and numeric_count >= 3)
                    or numeric_count >= 5
                    or "%" in compact
                ):
                    return False

        return True

    @staticmethod
    def _is_duplicate_image_note(
        note_text: str,
        slide_text_snapshot: str,
        existing_notes: list[str],
    ) -> bool:
        compact_note = re.sub(r"\s+", " ", str(note_text or "")).strip().lower()
        if not compact_note:
            return True

        compact_slide = re.sub(r"\s+", " ", str(slide_text_snapshot or "")).strip().lower()
        if compact_note in compact_slide:
            return True

        for existing in existing_notes:
            normalized_existing = re.sub(r"\s+", " ", str(existing or "")).strip().lower()
            if not normalized_existing:
                continue
            if compact_note == normalized_existing:
                return True

        return False

    @classmethod
    def _extract_table_lines(cls, shape) -> list[str]:
        lines: list[str] = []
        try:
            for row in shape.table.rows:
                cells = [cls._normalize_slide_line(cell.text) for cell in row.cells]
                cells = [cell for cell in cells if cell]
                if cells:
                    lines.append(" | ".join(cells))
        except Exception:
            return []
        return lines

    @classmethod
    def _extract_chart_lines(cls, shape) -> list[str]:
        lines: list[str] = []
        chart = getattr(shape, "chart", None)
        if chart is None:
            return lines

        categories: list[str] = []
        try:
            categories = [
                cls._normalize_slide_line(str(category))
                for category in chart.plots[0].categories
            ]
        except Exception:
            categories = []

        for series_index, series in enumerate(getattr(chart, "series", []), start=1):
            series_name = cls._normalize_slide_line(str(getattr(series, "name", "") or f"Series {series_index}"))
            try:
                values = list(getattr(series, "values", []))
            except Exception:
                values = []

            pairs: list[str] = []
            for value_index, value in enumerate(values[:8]):
                value_text = cls._normalize_slide_line(str(value))
                if not value_text:
                    continue

                if value_index < len(categories) and categories[value_index]:
                    pairs.append(f"{categories[value_index]}: {value_text}")
                else:
                    pairs.append(f"Point {value_index + 1}: {value_text}")

            if pairs:
                lines.append(f"{series_name}: {', '.join(pairs)}")
            elif series_name:
                lines.append(series_name)

            if len(lines) >= 10:
                break

        return lines
