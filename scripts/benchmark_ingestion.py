from __future__ import annotations

import argparse
import json
import statistics
import sys
import time
from dataclasses import dataclass
from pathlib import Path
from tempfile import TemporaryDirectory

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject


PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from app.core.config import get_settings
from app.core.embedding_factory import build_embeddings
from app.repositories.faiss_vector_store_repository import FaissVectorStoreRepository
from app.services.document_ingestion_service import DocumentIngestionService
from app.services.document_loader_registry import DocumentLoaderRegistry
from app.services.document_loaders.docx_document_loader import DocxDocumentLoader
from app.services.document_loaders.markdown_document_loader import MarkdownDocumentLoader
from app.services.document_loaders.pdf_document_loader import PdfDocumentLoader
from app.services.document_loaders.pptx_document_loader import PptxDocumentLoader
from app.services.document_loaders.text_document_loader import TextDocumentLoader
from app.services.runtime_metrics import RuntimeMetrics
from app.services.text_chunking_service import TextChunkingService


@dataclass(frozen=True)
class BenchmarkCase:
    name: str
    target_seconds: float
    files: list[Path]


def _paragraph(seed: int) -> str:
    return (
        f"Section {seed}: enterprise retrieval augmented generation pipeline validates semantic chunking, "
        "metadata alignment, reranking quality, and context compression for production-grade document intelligence."
    )


def _write_txt(path: Path, paragraphs: int = 180) -> None:
    content = "\n\n".join(_paragraph(index + 1) for index in range(paragraphs))
    path.write_text(content, encoding="utf-8")


def _write_md(path: Path, sections: int = 60) -> None:
    blocks: list[str] = ["# Enterprise RAG Benchmark Document"]
    for index in range(1, sections + 1):
        blocks.append(f"## Topic {index}")
        blocks.append(_paragraph(index))
        blocks.append("- ingestion speed")
        blocks.append("- semantic precision")
        blocks.append("- production scalability")
        blocks.append("```python")
        blocks.append("def benchmark_signal(value: int) -> int:")
        blocks.append("    return value * 2")
        blocks.append("```")
    path.write_text("\n\n".join(blocks), encoding="utf-8")


def _write_docx(path: Path, sections: int = 120) -> None:
    document = DocxDocument()
    document.add_heading("Enterprise RAG Ingestion Benchmark", level=1)
    for index in range(1, sections + 1):
        document.add_heading(f"Section {index}", level=2)
        document.add_paragraph(_paragraph(index))
        document.add_paragraph(
            "This paragraph preserves section context and paragraph boundaries for quality retrieval."
        )
    document.save(path)


def _write_pptx(path: Path, slides: int = 40) -> None:
    presentation = Presentation()
    layout = presentation.slide_layouts[1]

    for index in range(1, slides + 1):
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = f"Slide {index}: Retrieval Strategy"
        body = slide.placeholders[1].text_frame
        body.clear()

        body.text = "Hybrid search combines vector semantics and lexical matching."
        for line in (
            "Metadata hints improve document-type targeting.",
            "Reranking promotes coherent section-level evidence.",
            "Context compression reduces token waste before generation.",
        ):
            paragraph = body.add_paragraph()
            paragraph.text = line

        notes = slide.notes_slide.notes_text_frame
        notes.text = (
            f"Slide note {index}: operational benchmark guidance for enterprise deployment."
        )

    presentation.save(path)


def _escape_pdf_text(value: str) -> str:
    return value.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def _write_pdf(path: Path, pages: int = 100) -> None:
    writer = PdfWriter()

    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_ref = writer._add_object(font)
    resources = DictionaryObject(
        {
            NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref}),
        }
    )

    for index in range(1, pages + 1):
        page = writer.add_blank_page(width=612, height=792)
        page[NameObject("/Resources")] = resources

        lines = [
            f"PDF page {index} benchmark for enterprise retrieval.",
            "Paragraph-aware chunking keeps reading flow and section relevance.",
            "Local semantic embeddings reduce indexing cost for production workloads.",
            "Hybrid retrieval and reranking improve answer precision.",
        ]

        commands: list[str] = ["BT", "/F1 11 Tf", "72 740 Td"]
        for line_index, line in enumerate(lines):
            escaped = _escape_pdf_text(line)
            if line_index == 0:
                commands.append(f"({escaped}) Tj")
            else:
                commands.append(f"0 -16 Td ({escaped}) Tj")
        commands.append("ET")

        stream = DecodedStreamObject()
        stream.set_data("\n".join(commands).encode("latin-1", errors="ignore"))
        page[NameObject("/Contents")] = writer._add_object(stream)

    with path.open("wb") as output_file:
        writer.write(output_file)


def _prepare_dataset(dataset_dir: Path) -> list[BenchmarkCase]:
    dataset_dir.mkdir(parents=True, exist_ok=True)

    txt_path = dataset_dir / "benchmark-large.txt"
    md_path = dataset_dir / "benchmark-large.md"
    docx_path = dataset_dir / "benchmark-large.docx"
    pptx_path = dataset_dir / "benchmark-40-slides.pptx"
    pdf_path = dataset_dir / "benchmark-100-pages.pdf"

    _write_txt(txt_path)
    _write_md(md_path)
    _write_docx(docx_path)
    _write_pptx(pptx_path, slides=40)
    _write_pdf(pdf_path, pages=100)

    return [
        BenchmarkCase(name="TXT/MD", target_seconds=2.0, files=[txt_path, md_path]),
        BenchmarkCase(name="DOCX", target_seconds=5.0, files=[docx_path]),
        BenchmarkCase(name="PPTX 40 slides", target_seconds=8.0, files=[pptx_path]),
        BenchmarkCase(name="PDF 100 pages", target_seconds=15.0, files=[pdf_path]),
    ]


def _build_ingestion_service(
    *,
    index_dir: Path,
    settings,
    embeddings,
) -> tuple[DocumentIngestionService, RuntimeMetrics]:
    runtime_metrics = RuntimeMetrics()
    vector_store_repository = FaissVectorStoreRepository(
        index_dir=index_dir,
        embeddings=embeddings,
        embedding_batch_size=settings.embedding_batch_size,
        embedding_cache_enabled=settings.embedding_cache_enabled,
        runtime_metrics=runtime_metrics,
    )

    loader_registry = DocumentLoaderRegistry(
        loaders=[
            TextDocumentLoader(),
            MarkdownDocumentLoader(),
            DocxDocumentLoader(),
            PptxDocumentLoader(),
            PdfDocumentLoader(),
        ]
    )
    chunking_service = TextChunkingService(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
    )
    ingestion_service = DocumentIngestionService(
        loader_registry=loader_registry,
        chunking_service=chunking_service,
        vector_store_repository=vector_store_repository,
        max_file_workers=max(4, settings.ingestion_max_file_workers),
        runtime_metrics=runtime_metrics,
    )
    return ingestion_service, runtime_metrics


def _warm_up_embeddings(settings, embeddings, warmup_dir: Path) -> None:
    warmup_file = warmup_dir / "warmup.txt"
    warmup_file.parent.mkdir(parents=True, exist_ok=True)
    warmup_file.write_text("warmup semantic embedding call", encoding="utf-8")

    with TemporaryDirectory(prefix="aichatbox_warmup_") as temp_dir:
        index_dir = Path(temp_dir) / "faiss"
        service, _ = _build_ingestion_service(
            index_dir=index_dir,
            settings=settings,
            embeddings=embeddings,
        )
        service.ingest([warmup_file], metadata={"owner": "benchmark", "chat_id": "warmup"})


def _run_case(case: BenchmarkCase, *, settings, embeddings, runs: int) -> dict:
    durations: list[float] = []
    chunk_counts: list[int] = []
    metric_samples: list[dict] = []

    for run_index in range(1, runs + 1):
        with TemporaryDirectory(prefix="aichatbox_ingest_bench_") as temp_dir:
            index_dir = Path(temp_dir) / "faiss"
            service, runtime_metrics = _build_ingestion_service(
                index_dir=index_dir,
                settings=settings,
                embeddings=embeddings,
            )

            started_at = time.perf_counter()
            result = service.ingest(
                case.files,
                metadata={"owner": "benchmark", "chat_id": f"{case.name}-run-{run_index}"},
            )
            elapsed_seconds = time.perf_counter() - started_at

            durations.append(elapsed_seconds)
            chunk_counts.append(result.chunks_indexed)
            metric_samples.append(runtime_metrics.snapshot())

    average_seconds = statistics.mean(durations)
    min_seconds = min(durations)
    max_seconds = max(durations)
    average_chunks = int(round(statistics.mean(chunk_counts)))

    return {
        "name": case.name,
        "target_seconds": case.target_seconds,
        "runs": runs,
        "durations_seconds": [round(value, 3) for value in durations],
        "avg_seconds": round(average_seconds, 3),
        "min_seconds": round(min_seconds, 3),
        "max_seconds": round(max_seconds, 3),
        "avg_chunks_indexed": average_chunks,
        "status": "PASS" if average_seconds <= case.target_seconds else "FAIL",
        "latest_metrics": metric_samples[-1],
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="Benchmark enterprise ingestion pipeline")
    parser.add_argument("--runs", type=int, default=1, help="Number of benchmark runs per case")
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("tmp") / "benchmark_ingestion_results.json",
        help="Output JSON file path",
    )
    parser.add_argument(
        "--dataset-dir",
        type=Path,
        default=Path("tmp") / "benchmark_dataset",
        help="Directory for generated benchmark files",
    )
    parser.add_argument(
        "--force-local-embeddings",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Force local semantic embeddings during benchmark",
    )
    parser.add_argument(
        "--local-model",
        type=str,
        default="",
        help=(
            "Optional local embedding model override "
            "(for example: sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2)"
        ),
    )
    args = parser.parse_args()

    if args.runs < 1:
        raise ValueError("--runs must be >= 1")

    settings = get_settings().model_copy(deep=True)
    if args.force_local_embeddings:
        settings.local_semantic_embeddings_enabled = True
        settings.local_semantic_embeddings = True
    if args.local_model.strip():
        settings.local_embedding_model = args.local_model.strip()

    embeddings = build_embeddings(settings)
    embedding_provider = embeddings.__class__.__name__
    embedding_model = str(getattr(embeddings, "model_name", "unknown"))
    embedding_device = str(getattr(embeddings, "device", "unknown"))

    benchmark_cases = _prepare_dataset(args.dataset_dir)
    _warm_up_embeddings(settings, embeddings, args.dataset_dir)

    results = {
        "provider": embedding_provider,
        "model": embedding_model,
        "device": embedding_device,
        "chunk_size": settings.chunk_size,
        "chunk_overlap": settings.chunk_overlap,
        "embedding_batch_size": settings.embedding_batch_size,
        "cases": [],
    }

    print("Ingestion Benchmark (warm start)")
    print(f"Embedding provider: {embedding_provider}")
    print(f"Embedding model: {embedding_model}")
    print(f"Embedding device: {embedding_device}")

    for case in benchmark_cases:
        case_result = _run_case(case, settings=settings, embeddings=embeddings, runs=args.runs)
        results["cases"].append(case_result)
        print(
            f"- {case_result['name']}: avg={case_result['avg_seconds']}s "
            f"target<{case_result['target_seconds']}s status={case_result['status']} "
            f"chunks={case_result['avg_chunks_indexed']}"
        )

    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.output.write_text(json.dumps(results, ensure_ascii=True, indent=2), encoding="utf-8")
    print(f"Saved report: {args.output}")


if __name__ == "__main__":
    main()
